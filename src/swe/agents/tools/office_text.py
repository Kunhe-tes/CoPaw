# -*- coding: utf-8 -*-
"""Text extraction for binary Office document formats used by read_file.

read_file remains a text-first tool; when the target file is a binary Office
document (docx/xlsx/pptx, or legacy .doc via antiword), the raw text is
extracted here and then flows through the same line-range/truncation logic
as plain text files.
"""

from __future__ import annotations

import shutil
import subprocess
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from ..tool_failure import ToolExecutionError

_ZIP_MAGIC = b"PK\x03\x04"
_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
_PDF_MAGIC = b"%PDF-"

_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

# Mirror the text-file read cap (1GB) for extracted document text.
_MAX_EXTRACTED_CHARS = 1024 * 1024 * 1024

_DOCX_EXTS = {".docx", ".docm"}
_XLSX_EXTS = {".xlsx", ".xlsm", ".xltx", ".xltm"}
_PPTX_EXTS = {".pptx", ".pptm"}


def _unsupported(
    file_name: str, format_name: str, hint: str
) -> ToolExecutionError:
    """Build the canonical error for a recognized-but-unsupported format."""
    return ToolExecutionError(
        error_type="invalid_arguments",
        detail=f"Error: {file_name}: {format_name} is not supported by read_file. {hint}",
    )


def _zip_contains(file_path: str, prefix: str) -> bool:
    """Return True when any ZIP member name starts with prefix."""
    try:
        with zipfile.ZipFile(file_path) as zf:
            return any(name.startswith(prefix) for name in zf.namelist())
    except (zipfile.BadZipFile, OSError):
        return False


def detect_document_kind(file_path: str) -> str | None:
    """Detect the binary document format of a file by magic bytes.

    Returns a kind string ("docx", "xlsx", "pptx", "doc", "xls", "ppt",
    "zip", "ole2", "pdf") for recognized binary formats, or None when the
    file looks like plain text (the caller falls back to text reading).
    """
    try:
        with open(file_path, "rb") as fh:
            head = fh.read(8)
    except OSError:
        return None

    ext = Path(file_path).suffix.lower()

    if head.startswith(_ZIP_MAGIC):
        if ext in _DOCX_EXTS or _zip_contains(file_path, "word/"):
            return "docx"
        if ext in _XLSX_EXTS or _zip_contains(file_path, "xl/"):
            return "xlsx"
        if ext in _PPTX_EXTS or _zip_contains(file_path, "ppt/"):
            return "pptx"
        return "zip"

    if head.startswith(_OLE2_MAGIC):
        if ext == ".doc":
            return "doc"
        if ext == ".xls":
            return "xls"
        if ext == ".ppt":
            return "ppt"
        return "ole2"

    if head.startswith(_PDF_MAGIC):
        return "pdf"

    return None


def _docx_paragraph_text(p_elem: ET.Element) -> str:
    """Collect run text from a w:p paragraph, honoring tabs and breaks."""
    parts: list[str] = []
    for node in p_elem.iter():
        tag = node.tag
        if tag == _W + "t":
            parts.append(node.text or "")
        elif tag == _W + "tab":
            parts.append("\t")
        elif tag in (_W + "br", _W + "cr"):
            parts.append("\n")
    return "".join(parts)


def _docx_table_lines(tbl_elem: ET.Element) -> list[str]:
    """Render a w:tbl as one line per row with cells joined by ' | '."""
    lines: list[str] = []
    for tr in tbl_elem.findall(_W + "tr"):
        cells: list[str] = []
        for tc in tr.findall(_W + "tc"):
            cell_parts: list[str] = []
            for p in tc.findall(_W + "p"):
                text = _docx_paragraph_text(p)
                if text.strip():
                    cell_parts.append(text)
            cells.append(" ".join(cell_parts))
        line = " | ".join(cells).rstrip(" |")
        if line:
            lines.append(line)
    return lines


def _extract_docx_text(file_path: str) -> str:
    """Extract plain text from the body of a .docx/.docm file."""
    name = Path(file_path).name
    try:
        with zipfile.ZipFile(file_path) as zf:
            if "word/document.xml" not in zf.namelist():
                raise _unsupported(
                    name,
                    "not a valid Word document (missing word/document.xml)",
                    "Re-save it as .docx and retry.",
                )
            root = ET.fromstring(zf.read("word/document.xml"))
    except zipfile.BadZipFile as exc:
        raise _unsupported(
            name,
            "corrupt or encrypted Word document (bad ZIP archive)",
            "Re-save it without encryption and retry.",
        ) from exc
    except ET.ParseError as exc:
        raise _unsupported(
            name,
            "corrupt Word document (unparseable document.xml)",
            "Re-save it as .docx and retry.",
        ) from exc

    body = root.find(_W + "body")
    lines: list[str] = []
    total_chars = 0
    if body is not None:
        for child in body:
            if child.tag == _W + "p":
                text = _docx_paragraph_text(child)
                if not text.strip():
                    continue
            elif child.tag == _W + "tbl":
                table_lines = _docx_table_lines(child)
                for line in table_lines:
                    if total_chars + len(line) > _MAX_EXTRACTED_CHARS:
                        raise _unsupported(
                            name,
                            "extracted text exceeds the size limit",
                            "Extract the content with pandoc or the docx skill instead.",
                        )
                    lines.append(line)
                    total_chars += len(line)
                continue
            else:
                continue
            if total_chars + len(text) > _MAX_EXTRACTED_CHARS:
                raise _unsupported(
                    name,
                    "extracted text exceeds the size limit",
                    "Extract the content with pandoc or the docx skill instead.",
                )
            lines.append(text)
            total_chars += len(text)
    return "\n".join(lines)


def _extract_xlsx_text(file_path: str) -> str:
    """Extract cell text from every sheet of a .xlsx/.xlsm file."""
    name = Path(file_path).name
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise _unsupported(
            name,
            "Excel extraction requires the 'openpyxl' package",
            "Install openpyxl or use the xlsx skill instead.",
        ) from exc

    wb = None
    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        lines: list[str] = []
        total_chars = 0
        for ws in wb.worksheets:
            header = f"=== Sheet: {ws.title} ==="
            if total_chars + len(header) > _MAX_EXTRACTED_CHARS:
                break
            lines.append(header)
            total_chars += len(header)
            for row in ws.iter_rows(values_only=True):
                values = ["" if v is None else str(v) for v in row]
                if not any(values):
                    continue
                line = " | ".join(values).rstrip(" |")
                if total_chars + len(line) > _MAX_EXTRACTED_CHARS:
                    break
                lines.append(line)
                total_chars += len(line)
        return "\n".join(lines)
    except Exception as exc:
        if isinstance(exc, ToolExecutionError):
            raise
        raise _unsupported(
            name,
            "could not read Excel content",
            "Re-save it as .xlsx (or a CSV) and retry.",
        ) from exc
    finally:
        if wb is not None:
            wb.close()


def _collect_pptx_shape_text(shape: object, lines: list[str]) -> None:
    """Collect text from a pptx shape, recursing into groups."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        MSO_SHAPE_TYPE = None  # type: ignore[assignment]

    if (
        MSO_SHAPE_TYPE is not None
        and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
    ):
        for sub in getattr(shape, "shapes", []):
            _collect_pptx_shape_text(sub, lines)
        return

    if getattr(shape, "has_text_frame", False):
        for para in getattr(shape.text_frame, "paragraphs", []):
            text = getattr(para, "text", "")
            if text and text.strip():
                lines.append(text.strip())

    if getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            cells = [
                cell.text.replace("\n", " ").strip() for cell in row.cells
            ]
            if any(cells):
                lines.append(" | ".join(cells).rstrip(" |"))


def _extract_pptx_text(file_path: str) -> str:
    """Extract text from every slide of a .pptx/.pptm file."""
    name = Path(file_path).name
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise _unsupported(
            name,
            "PowerPoint extraction requires the 'python-pptx' package",
            "Install python-pptx or use the pptx skill instead.",
        ) from exc

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        raise _unsupported(
            name,
            "could not read PowerPoint content",
            "Re-save it as .pptx and retry.",
        ) from exc

    lines: list[str] = []
    total_chars = 0
    over_cap = False
    for index, slide in enumerate(prs.slides, start=1):
        if over_cap:
            break
        header = f"=== Slide {index} ==="
        if total_chars + len(header) > _MAX_EXTRACTED_CHARS:
            break
        lines.append(header)
        total_chars += len(header)
        for shape in slide.shapes:
            before = len(lines)
            _collect_pptx_shape_text(shape, lines)
            for line in lines[before:]:
                if total_chars + len(line) > _MAX_EXTRACTED_CHARS:
                    del lines[before:]
                    over_cap = True
                    break
                total_chars += len(line)
            if over_cap:
                break
    return "\n".join(lines)


def _extract_doc_text(file_path: str) -> str:
    """Extract text from a legacy binary .doc file via antiword."""
    name = Path(file_path).name
    exe = shutil.which("antiword")
    if not exe:
        raise _unsupported(
            name,
            "legacy .doc files require the 'antiword' tool",
            "Convert the file to .docx (e.g. with LibreOffice) and retry.",
        )
    try:
        # No check=: a non-zero exit is a formatted error, not an exception.
        proc = subprocess.run(  # pylint: disable=subprocess-run-check
            [exe, file_path],
            capture_output=True,
            timeout=120,
        )
    except (subprocess.TimeoutExpired, OSError) as exc:
        raise _unsupported(
            name,
            "antiword failed to extract the document",
            "Convert the file to .docx (e.g. with LibreOffice) and retry.",
        ) from exc
    if proc.returncode != 0:
        stderr = proc.stderr.decode("utf-8", errors="replace").strip()
        raise _unsupported(
            name,
            "antiword could not decode this .doc file",
            f"antiword: {stderr or f'exit code {proc.returncode}'}. "
            "Convert the file to .docx and retry.",
        )
    return proc.stdout.decode("utf-8", errors="replace")


def extract_document_text(file_path: str, kind: str) -> str:
    """Extract plain text for a detected binary document kind."""
    name = Path(file_path).name
    if kind == "docx":
        return _extract_docx_text(file_path)
    if kind == "xlsx":
        return _extract_xlsx_text(file_path)
    if kind == "pptx":
        return _extract_pptx_text(file_path)
    if kind == "doc":
        return _extract_doc_text(file_path)
    if kind == "xls":
        raise _unsupported(
            name,
            "legacy .xls format is not supported",
            "Convert the file to .xlsx (e.g. with LibreOffice) and retry.",
        )
    if kind == "ppt":
        raise _unsupported(
            name,
            "legacy .ppt format is not supported",
            "Convert the file to .pptx (e.g. with LibreOffice) and retry.",
        )
    if kind == "ole2":
        raise _unsupported(
            name,
            "OLE2 compound document is not supported",
            "Convert it to a modern Office format and retry.",
        )
    if kind == "pdf":
        raise _unsupported(
            name,
            "PDF files are not supported",
            "Convert the PDF to text or use dedicated PDF tooling.",
        )
    if kind == "zip":
        raise _unsupported(
            name,
            "ZIP archive is not a supported document format",
            "Only docx/xlsx/pptx archives can be read as text.",
        )
    raise _unsupported(name, f"unknown format {kind!r}", "")
