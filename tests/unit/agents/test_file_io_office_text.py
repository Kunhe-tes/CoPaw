# -*- coding: utf-8 -*-
from __future__ import annotations

from pathlib import Path
import zipfile

import pytest

from swe.agents.tool_failure import ToolExecutionError
from swe.agents.tools import file_io
from swe.agents.tools.office_text import (
    detect_document_kind,
    extract_document_text,
)

_DOCX_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"


def _xml_escape(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _make_docx(
    path: Path,
    paragraphs: list[str],
    table_rows: list[list[str]] | None = None,
) -> None:
    body_parts: list[str] = []
    for para in paragraphs:
        body_parts.append(
            '<w:p><w:r><w:t xml:space="preserve">'
            f"{_xml_escape(para)}</w:t></w:r></w:p>"
        )
    if table_rows:
        row_xml: list[str] = []
        for row in table_rows:
            cells = "".join(
                '<w:tc><w:p><w:r><w:t xml:space="preserve">'
                f"{_xml_escape(cell)}</w:t></w:r></w:p></w:tc>"
                for cell in row
            )
            row_xml.append(f"<w:tr>{cells}</w:tr>")
        body_parts.append(f"<w:tbl>{''.join(row_xml)}</w:tbl>")
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:document xmlns:w="{_DOCX_NS}">'
        f"<w:body>{''.join(body_parts)}<w:sectPr/></w:body></w:document>"
    )
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0"?><Types '
            'xmlns="http://schemas.openxmlformats.org/package/2006/content-types"/>',
        )
        zf.writestr("word/document.xml", document)


def _make_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Value"])
    ws.append(["Alice", 42])
    wb.save(str(path))


def _make_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 914400, 457200)
    box.text_frame.text = "Hello Slides"
    prs.save(str(path))


def test_detect_text_file_returns_none(tmp_path: Path) -> None:
    target = tmp_path / "note.md"
    target.write_text("alpha\nbeta", encoding="utf-8")
    assert detect_document_kind(str(target)) is None


def test_detect_docx_by_extension(tmp_path: Path) -> None:
    target = tmp_path / "report.docx"
    _make_docx(target, ["Hello"])
    assert detect_document_kind(str(target)) == "docx"


def test_detect_docx_by_content_without_extension(tmp_path: Path) -> None:
    target = tmp_path / "report.bin"
    _make_docx(target, ["Hello"])
    assert detect_document_kind(str(target)) == "docx"


def test_detect_xlsx_and_pptx(tmp_path: Path) -> None:
    xlsx = tmp_path / "data.xlsx"
    _make_xlsx(xlsx)
    assert detect_document_kind(str(xlsx)) == "xlsx"

    pptx = tmp_path / "deck.pptx"
    _make_pptx(pptx)
    assert detect_document_kind(str(pptx)) == "pptx"


def test_detect_legacy_ole_and_pdf(tmp_path: Path) -> None:
    doc = tmp_path / "old.doc"
    doc.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 16)
    assert detect_document_kind(str(doc)) == "doc"

    xls = tmp_path / "old.xls"
    xls.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 16)
    assert detect_document_kind(str(xls)) == "xls"

    pdf = tmp_path / "doc.pdf"
    pdf.write_bytes(b"%PDF-1.4\n1 0 obj\n")
    assert detect_document_kind(str(pdf)) == "pdf"


def test_extract_docx_paragraphs_and_table(tmp_path: Path) -> None:
    target = tmp_path / "report.docx"
    _make_docx(
        target,
        ["Hello", "World"],
        table_rows=[["A", "B"], ["C", "D"]],
    )
    text = extract_document_text(str(target), "docx")
    assert text == "Hello\nWorld\nA | B\nC | D"


def test_extract_xlsx_sheets_and_rows(tmp_path: Path) -> None:
    target = tmp_path / "data.xlsx"
    _make_xlsx(target)
    text = extract_document_text(str(target), "xlsx")
    assert "=== Sheet: Data ===" in text
    assert "Name | Value" in text
    assert "Alice | 42" in text


def test_extract_pptx_slides(tmp_path: Path) -> None:
    target = tmp_path / "deck.pptx"
    _make_pptx(target)
    text = extract_document_text(str(target), "pptx")
    assert "=== Slide 1 ===" in text
    assert "Hello Slides" in text


@pytest.mark.asyncio
async def test_read_file_docx_end_to_end(tmp_path: Path, monkeypatch) -> None:
    target = tmp_path / "report.docx"
    _make_docx(
        target,
        ["Hello", "World"],
        table_rows=[["A", "B"], ["C", "D"]],
    )
    monkeypatch.setattr(file_io, "_resolve_file_path", lambda _: str(target))

    result = await file_io.read_file("logical/report.docx")
    assert result.content[0]["text"] == "Hello\nWorld\nA | B\nC | D"

    ranged = await file_io.read_file(
        "logical/report.docx", start_line=2, end_line=3
    )
    ranged_text = ranged.content[0]["text"]
    assert ranged_text.startswith("World\nA | B")
    assert "start_line=4 to read more." in ranged_text


@pytest.mark.asyncio
async def test_read_file_pdf_raises_clear_error(
    tmp_path: Path, monkeypatch
) -> None:
    target = tmp_path / "doc.pdf"
    target.write_bytes(b"%PDF-1.4\n1 0 obj\n")
    monkeypatch.setattr(file_io, "_resolve_file_path", lambda _: str(target))

    with pytest.raises(ToolExecutionError) as exc_info:
        await file_io.read_file("logical/doc.pdf")

    assert exc_info.value.error_type == "invalid_arguments"
    assert "PDF files are not supported" in exc_info.value.detail


@pytest.mark.asyncio
async def test_read_file_legacy_xls_raises_clear_error(
    tmp_path: Path, monkeypatch
) -> None:
    target = tmp_path / "old.xls"
    target.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 16)
    monkeypatch.setattr(file_io, "_resolve_file_path", lambda _: str(target))

    with pytest.raises(ToolExecutionError) as exc_info:
        await file_io.read_file("logical/old.xls")

    assert exc_info.value.error_type == "invalid_arguments"
    assert "legacy .xls format is not supported" in exc_info.value.detail
